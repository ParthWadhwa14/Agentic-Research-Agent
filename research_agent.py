import base64
import io
import json
import mimetypes
import operator
import os
import re
import tempfile
from datetime import datetime
from pathlib import Path
from zoneinfo import ZoneInfo
from typing import Annotated, Any, Dict, List, Optional, Sequence, TypedDict, Union

from dotenv import load_dotenv
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_groq import ChatGroq
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.utilities import GoogleSerperAPIWrapper, WikipediaAPIWrapper
from langchain_community.tools import WikipediaQueryRun
from langgraph.graph import END, START, StateGraph

from memory_system import ContextBuilder, InMemoryMemoryStore, MemoryRecord, MemoryStore
from media_search import collect_slide_images
from data_analysis_agent import analyze_data_file as run_local_data_analysis

load_dotenv()

# ==========================================
# 1. STATE SCHEMA
# ==========================================
class ResearchState(TypedDict):
    query: str
    format_preference: str
    memory_context: str
    task_intent: str
    input_context: str
    source_files: List[str]
    chart_paths: List[str]
    pptx_path: str
    plan: str
    research_data: Annotated[str, operator.add]
    current_gaps: str
    iteration: int
    final_report: str


# ==========================================
# 2. INITIALIZE MODELS & TOOLS
# ==========================================
llm = ChatGroq(model="meta-llama/llama-4-scout-17b-16e-instruct", temperature=0.2)
llm_vision = ChatGroq(model="meta-llama/llama-4-scout-17b-16e-instruct", temperature=0.2)
web_search_tool = GoogleSerperAPIWrapper()
wiki_search_tool = WikipediaQueryRun(api_wrapper=WikipediaAPIWrapper())

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
DATA_EXTENSIONS = {".csv", ".xlsx", ".xls"}

default_memory_store = InMemoryMemoryStore()


def build_current_datetime_context() -> str:
    """Provides a reliable clock anchor for relative dates in research mode."""
    local_time = datetime.now().astimezone()
    utc_time = datetime.now(ZoneInfo("UTC"))
    india_time = datetime.now(ZoneInfo("Asia/Kolkata"))
    return "\n".join([
        f"Current server time: {local_time.strftime('%Y-%m-%d %I:%M %p %Z')}",
        f"Current UTC time: {utc_time.strftime('%Y-%m-%d %H:%M:%S %Z')}",
        f"Current India time: {india_time.strftime('%Y-%m-%d %I:%M %p %Z')} (Asia/Kolkata)",
        f"Current year: {india_time.year}",
    ])


def clean_text(value: Any, max_chars: int = 700) -> str:
    if value is None:
        return ""
    text = re.sub(r"\s+", " ", str(value)).strip()
    return text[:max_chars].rstrip()


def format_serper_results(payload: Dict[str, Any], query: str, max_items: int = 7) -> str:
    lines = [f"Search query: {query}"]
    direct_result = clean_text(payload.get("answerBox") or payload.get("knowledgeGraph"), 1000)
    if direct_result:
        lines.append(f"Direct result: {direct_result}")

    results = payload.get("organic") or payload.get("news") or payload.get("places") or []
    for index, item in enumerate(results[:max_items], start=1):
        if not isinstance(item, dict):
            continue
        title = clean_text(item.get("title"), 180)
        snippet = clean_text(item.get("snippet") or item.get("description"), 550)
        link = clean_text(item.get("link"), 260)
        date = clean_text(item.get("date"), 80)
        source = clean_text(item.get("source"), 120)
        meta = " · ".join(part for part in [source, date] if part)
        meta_text = f" ({meta})" if meta else ""
        if title or snippet or link:
            lines.append(f"{index}. {title}{meta_text}\n   {snippet}\n   URL: {link}")
    return "\n".join(lines)


def serper_search(query: str) -> str:
    if not os.getenv("SERPER_API_KEY"):
        return "Web search unavailable: SERPER_API_KEY is not set in the backend environment."
    try:
        payload = GoogleSerperAPIWrapper(k=8, gl="us", hl="en").results(query)
        return format_serper_results(payload, query)
    except Exception as exc:
        return f"Search failed for '{query}': {exc}"


def search_memory(query: str, limit: int = 5) -> str:
    """Tool-style helper for retrieving only relevant long-term memories."""
    results = default_memory_store.search(query, limit=limit)
    if not results:
        return "No relevant memories found."
    return "\n".join(
        f"- {item.memory.title} ({item.memory.memory_type}, score={item.score:.2f}): {item.memory.summary}"
        for item in results
    )


def get_report(report_id: str) -> str:
    """Tool-style helper that returns metadata/path for a stored report without injecting full contents."""
    for memory in default_memory_store.memories:
        if memory.id == report_id and memory.memory_type == "research_report":
            return json.dumps({
                "id": memory.id,
                "title": memory.title,
                "summary": memory.summary,
                "tags": memory.tags,
                "file_path": memory.file_path,
                "created_at": memory.created_at.isoformat(),
            })
    return f"No report metadata found for id: {report_id}"


def get_user_preference(key: str) -> str:
    """Tool-style helper for preference lookup."""
    for memory in default_memory_store.memories:
        if memory.memory_type == "user_preference" and (key in memory.tags or key.lower() in memory.title.lower()):
            return memory.summary
    return f"No user preference found for key: {key}"


# ==========================================
# 3. INPUT, MATH, DATA, AND PPT HELPERS
# ==========================================
def _as_list(paths: Optional[Union[str, Sequence[str]]]) -> List[str]:
    if not paths:
        return []
    if isinstance(paths, str):
        return [paths]
    return [str(path) for path in paths if path]


def extract_pdf_text(pdf_path: str, max_chars: int = 12000) -> str:
    """Extracts text from a PDF for grounding research."""
    if not os.path.exists(pdf_path):
        return f"[PDF missing: {pdf_path}]"
    try:
        pages = PyPDFLoader(pdf_path).load()
        text = "\n\n".join(page.page_content for page in pages)
        return text[:max_chars]
    except Exception as exc:
        return f"[PDF extraction failed for {pdf_path}: {exc}]"


def encode_image(image_path: str) -> str:
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def prepare_image_data_url(image_path: str, max_side: int = 1600, jpeg_quality: int = 85) -> str:
    """Builds a vision-model-compatible image data URL, resizing if Pillow is available."""
    mime_type = mimetypes.guess_type(image_path)[0] or "image/jpeg"
    try:
        from PIL import Image

        with Image.open(image_path) as image:
            image = image.convert("RGB")
            image.thumbnail((max_side, max_side))
            buffer = io.BytesIO()
            image.save(buffer, format="JPEG", quality=jpeg_quality, optimize=True)
            encoded = base64.b64encode(buffer.getvalue()).decode("utf-8")
            return f"data:image/jpeg;base64,{encoded}"
    except Exception:
        return f"data:{mime_type};base64,{encode_image(image_path)}"


def analyze_image(image_path: str) -> str:
    """Uses the vision model to summarize visual evidence and OCR text."""
    if not os.path.exists(image_path):
        return f"[Image missing: {image_path}]"

    image_size = os.path.getsize(image_path)
    image_mime = mimetypes.guess_type(image_path)[0] or "unknown"
    prompt = (
        f"Analyze this uploaded image for a research report.\n"
        f"Upload metadata: MIME/type={image_mime}; size={image_size} bytes.\n"
        "Return: overall description, "
        "important visible objects/data, exact OCR text, chart/table interpretation if present, "
        "and any limitations or uncertainty."
    )
    message = HumanMessage(
        content=[
            {"type": "text", "text": prompt},
            {
                "type": "image_url",
                "image_url": {"url": prepare_image_data_url(image_path)},
            },
        ]
    )
    try:
        return llm_vision.invoke([message]).content
    except Exception as exc:
        return f"[Image analysis failed for {image_path}: {exc}]"


def solve_math_problem(query: str) -> str:
    """Math agent for algebra, calculus, statistics, optimization, and numerical methods."""
    prompt = f"""
    You are a Math Agent with access conceptually to SymPy, NumPy, SciPy, and a Python sandbox.
    Determine whether the user query contains a mathematical or quantitative problem.

    If it does not, output exactly: NO_MATH_REQUIRED

    If it does, solve it using this workflow:
    1. Understand the problem.
    2. Convert it to symbolic/math form.
    3. Solve using appropriate SymPy/NumPy/SciPy-style reasoning.
    4. Verify the result.
    5. Explain the steps clearly.

    User query:
    {query}
    """
    try:
        response = llm.invoke([HumanMessage(content=prompt)]).content
        return "" if "NO_MATH_REQUIRED" in response else response
    except Exception as exc:
        return f"Math agent failed: {exc}"


def analyze_data_file(file_path: str, output_dir: Optional[str] = None) -> Dict[str, Any]:
    """Runs lightweight data analysis on CSV/Excel files and emits chart paths."""
    result = {"summary": "", "chart_paths": []}
    if not os.path.exists(file_path):
        result["summary"] = f"[Data file missing: {file_path}]"
        return result

    try:
        import pandas as pd
        import matplotlib.pyplot as plt
    except Exception as exc:
        result["summary"] = f"[Data analysis dependencies missing: {exc}]"
        return result

    try:
        suffix = Path(file_path).suffix.lower()
        df = pd.read_csv(file_path) if suffix == ".csv" else pd.read_excel(file_path)
        output = Path(output_dir or tempfile.mkdtemp(prefix="research_charts_"))
        output.mkdir(parents=True, exist_ok=True)

        missing = df.isna().sum().to_dict()
        numeric_df = df.select_dtypes(include="number")
        summary_parts = [
            f"Data file: {file_path}",
            f"Shape: {df.shape[0]} rows x {df.shape[1]} columns",
            f"Columns: {', '.join(map(str, df.columns))}",
            f"Missing values by column: {missing}",
        ]

        if not numeric_df.empty:
            summary_parts.append("Summary statistics:\n" + numeric_df.describe().to_markdown())
            correlation = numeric_df.corr(numeric_only=True)
            summary_parts.append("Correlation matrix:\n" + correlation.to_markdown())

            ax = numeric_df.hist(figsize=(10, 7), bins=20)
            plt.tight_layout()
            hist_path = output / f"{Path(file_path).stem}_histograms.png"
            plt.savefig(hist_path)
            plt.close("all")
            result["chart_paths"].append(str(hist_path))

            if len(numeric_df.columns) >= 2:
                plt.figure(figsize=(8, 6))
                plt.imshow(correlation, cmap="coolwarm", vmin=-1, vmax=1)
                plt.colorbar(label="Correlation")
                plt.xticks(range(len(correlation.columns)), correlation.columns, rotation=45, ha="right")
                plt.yticks(range(len(correlation.columns)), correlation.columns)
                plt.title("Correlation Heatmap")
                plt.tight_layout()
                corr_path = output / f"{Path(file_path).stem}_correlation.png"
                plt.savefig(corr_path)
                plt.close("all")
                result["chart_paths"].append(str(corr_path))

            outlier_notes = []
            for column in numeric_df.columns:
                q1 = numeric_df[column].quantile(0.25)
                q3 = numeric_df[column].quantile(0.75)
                iqr = q3 - q1
                outliers = numeric_df[(numeric_df[column] < q1 - 1.5 * iqr) | (numeric_df[column] > q3 + 1.5 * iqr)]
                if not outliers.empty:
                    outlier_notes.append(f"{column}: {len(outliers)} potential outliers")
            summary_parts.append("Outlier detection: " + ("; ".join(outlier_notes) if outlier_notes else "No strong IQR outliers found."))
        else:
            summary_parts.append("No numeric columns detected for statistical charts/correlation.")

        result["summary"] = "\n\n".join(summary_parts)
        return result
    except Exception as exc:
        result["summary"] = f"[Data analysis failed for {file_path}: {exc}]"
        return result


def _extract_json_array(text: str) -> Optional[List[Any]]:
    """Extracts the first top-level JSON array from an LLM response."""
    cleaned = text.replace("```json", "").replace("```", "").strip()
    try:
        parsed = json.loads(cleaned)
        return parsed if isinstance(parsed, list) else None
    except json.JSONDecodeError:
        pass

    start = cleaned.find("[")
    end = cleaned.rfind("]")
    if start == -1 or end == -1 or end <= start:
        return None
    try:
        parsed = json.loads(cleaned[start:end + 1])
        return parsed if isinstance(parsed, list) else None
    except json.JSONDecodeError:
        return None


def _clean_slide_text(text: Any, max_chars: int = 180) -> str:
    text = re.sub(r"\s+", " ", str(text or "")).strip(" -•*\t")
    text = re.sub(r"^\[(?:source|citation|reference)\]\s*", "", text, flags=re.IGNORECASE)
    return text[:max_chars].rstrip()


def _dedupe_preserve_order(items: Sequence[str]) -> List[str]:
    seen = set()
    deduped = []
    for item in items:
        key = re.sub(r"\W+", "", item.lower())
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(item)
    return deduped


def _derive_deck_title(query: str, report: str) -> str:
    for line in report.splitlines():
        title = re.sub(r"^#{1,6}\s*", "", line).strip()
        if 12 <= len(title) <= 110 and not title.lower().startswith(("executive summary", "direct answer")):
            return title
    return query[:110] or "Research Presentation"


def _extract_references(report: str, max_items: int = 6) -> List[str]:
    references = []
    url_pattern = re.compile(r"https?://[^\s)\]]+")
    for line in report.splitlines():
        if url_pattern.search(line) or re.match(r"^\s*(?:[-*]|\d+\.)?\s*(?:references?|sources?)\b", line, re.IGNORECASE):
            clean_line = _clean_slide_text(line, 220)
            if clean_line and not clean_line.lower() in {"references", "sources"}:
                references.append(clean_line)
    return _dedupe_preserve_order(references)[:max_items]


def presentation_planner_agent(report: str, query: str) -> List[Dict[str, Any]]:
    """Plans a dynamic slide deck from either a chat answer or a research report."""
    recommended_sections = [
        "Opening",
        "Context",
        "Method",
        "Key Findings",
        "Evidence",
        "Comparison",
        "Risks",
        "Recommendations",
        "Next Steps",
        "References",
    ]
    prompt = f"""
    You are a senior presentation strategist.
    Convert the source material into a boardroom-quality slide deck for the user's request.

    Goal:
    - Create a clear narrative arc, not a copied report.
    - Use strong slide titles that communicate the insight, not generic labels.
    - Produce 6-10 slides unless the material is genuinely tiny.
    - Prefer 3-5 short, punchy bullets per content slide.
    - Keep each bullet under 18 words.
    - Add one "takeaway" per important slide.
    - Add speaker notes with useful context, evidence, caveats, and transitions; do not just repeat bullets.
    - Split dense topics and merge weak/repetitive sections.
    - Include a recommendations or next-steps slide when the query asks for decisions, strategy, comparison, or advice.
    - Include a references slide if source URLs or citations exist.
    - Include charts/images only when they improve understanding.

    Recommended sections to consider, not mandatory and not fixed:
    {recommended_sections}

    Return valid JSON only, with this shape:
    [
      {{
        "title": "Slide title",
        "section": "short section label",
        "slide_type": "title | agenda | insight | comparison | data | recommendation | risk | reference",
        "takeaway": "one sentence key takeaway",
        "bullets": ["short bullet", "short bullet"],
        "visual": "none | chart | image",
        "speaker_notes": "brief speaker notes"
      }}
    ]

    Query:
    {query}

    Report:
    {report[:18000]}
    """
    try:
        response = llm.invoke([HumanMessage(content=prompt)]).content.strip()
        parsed = _extract_json_array(response)
        if isinstance(parsed, list):
            slides = []
            for item in parsed:
                if not isinstance(item, dict):
                    continue
                title = _clean_slide_text(item.get("title", ""), 95)
                bullets = _normalize_slide_bullets(item.get("bullets", []))
                visual = str(item.get("visual", "none")).strip().lower()
                if title:
                    slides.append({
                        "title": title,
                        "section": _clean_slide_text(item.get("section", "Insight"), 32) or "Insight",
                        "slide_type": _clean_slide_text(item.get("slide_type", "insight"), 24).lower() or "insight",
                        "takeaway": _clean_slide_text(item.get("takeaway", ""), 190),
                        "bullets": bullets[:5],
                        "visual": visual if visual in {"chart", "image"} else "none",
                        "speaker_notes": _clean_slide_text(item.get("speaker_notes", ""), 1200),
                    })
            if slides:
                return _polish_slide_plan(slides, query, report)
    except Exception:
        pass

    return _polish_slide_plan(_section_slides_from_report(report), query, report)


def _polish_slide_plan(slides: List[Dict[str, Any]], query: str, report: str, max_slides: int = 11) -> List[Dict[str, Any]]:
    """Normalizes LLM/fallback slide plans into a complete, presentable deck."""
    cleaned_slides: List[Dict[str, Any]] = []
    for slide in slides:
        title = _clean_slide_text(slide.get("title", ""), 95)
        if not title or title.lower() in {"title", "overview"}:
            continue
        bullets = _dedupe_preserve_order(_normalize_slide_bullets(slide.get("bullets", [])))[:5]
        if not bullets and slide.get("takeaway"):
            bullets = [_clean_slide_text(slide["takeaway"], 160)]
        cleaned_slides.append({
            "title": title,
            "section": _clean_slide_text(slide.get("section", "Insight"), 32) or "Insight",
            "slide_type": _clean_slide_text(slide.get("slide_type", "insight"), 24).lower() or "insight",
            "takeaway": _clean_slide_text(slide.get("takeaway", ""), 190),
            "bullets": bullets[:5] or ["See speaker notes for supporting detail."],
            "visual": slide.get("visual", "none") if slide.get("visual") in {"chart", "image"} else "none",
            "speaker_notes": _clean_slide_text(slide.get("speaker_notes", ""), 1200),
        })

    if not cleaned_slides:
        cleaned_slides = [{
            "title": _derive_deck_title(query, report),
            "section": "Overview",
            "slide_type": "insight",
            "takeaway": _clean_slide_text(report, 180),
            "bullets": _normalize_slide_bullets(report)[:5] or [_clean_slide_text(report, 160)],
            "visual": "none",
            "speaker_notes": _clean_slide_text(report, 900),
        }]

    needs_agenda = len(cleaned_slides) >= 4 and cleaned_slides[0].get("slide_type") != "agenda"
    if needs_agenda:
        agenda_items = _dedupe_preserve_order([slide["section"] for slide in cleaned_slides])[:5]
        cleaned_slides.insert(0, {
            "title": "Today’s Narrative",
            "section": "Opening",
            "slide_type": "agenda",
            "takeaway": "The deck moves from context to evidence to action.",
            "bullets": agenda_items or ["Context", "Findings", "Recommendations"],
            "visual": "none",
            "speaker_notes": "Use this slide to preview the story arc and set expectations for the audience.",
        })

    references = _extract_references(report)
    has_reference_slide = any(slide.get("slide_type") == "reference" or "reference" in slide["title"].lower() for slide in cleaned_slides)
    if references and not has_reference_slide:
        cleaned_slides.append({
            "title": "Selected Sources",
            "section": "References",
            "slide_type": "reference",
            "takeaway": "Claims should be checked against the linked source material.",
            "bullets": references,
            "visual": "none",
            "speaker_notes": "These are the highest-signal references detected in the generated report.",
        })

    return cleaned_slides[:max_slides]


def _section_slides_from_report(report: str, max_slides: int = 10) -> List[Dict[str, Any]]:
    """Fallback parser that creates usable slides from Markdown/plain text sections."""
    fallback_titles = [
        "Executive Snapshot",
        "Context That Matters",
        "Most Important Findings",
        "Evidence And Analysis",
        "Risks And Unknowns",
        "Recommended Actions",
        "Closing Takeaway",
    ]
    chunks = [chunk.strip() for chunk in re.split(r"\n(?=#{1,3}\s+|\d+\.\s+|Slide\s+\d+[:.)-])", report) if chunk.strip()]
    slides = []
    for index, chunk in enumerate(chunks[:max_slides] or [report[:1400]]):
        lines = [line.strip() for line in chunk.splitlines() if line.strip()]
        raw_title = re.sub(r"^(#{1,6}\s+|\d+\.\s+|Slide\s+\d+[:.)-]\s*)", "", lines[0] if lines else "").strip()
        title = raw_title if 8 <= len(raw_title) <= 95 else fallback_titles[index] if index < len(fallback_titles) else f"Supporting Detail {index + 1}"
        bullets = []
        for line in lines[1:]:
            stripped = _clean_slide_text(line, 170)
            if stripped and not stripped.startswith("#") and not re.match(r"^https?://", stripped):
                bullets.append(stripped)
        if not bullets:
            words = " ".join(lines[1:] or lines[:1]).split()
            bullets = [" ".join(words[item:item + 16]) for item in range(0, min(len(words), 80), 16)]
        takeaway = bullets[0] if bullets else _clean_slide_text(chunk, 180)
        slides.append({
            "title": title[:95],
            "section": "Summary",
            "slide_type": "insight",
            "takeaway": takeaway[:180],
            "bullets": bullets[:5] or [chunk[:240]],
            "visual": "none",
            "speaker_notes": _clean_slide_text(chunk, 900),
        })
    return slides


def _normalize_slide_bullets(content: Any) -> List[str]:
    if isinstance(content, list):
        return [_clean_slide_text(item, 175) for item in content if _clean_slide_text(item, 175)]
    if isinstance(content, str):
        parts = re.split(r"\n+|(?:^|\s)[-*•]\s+", content)
        return [_clean_slide_text(item, 175) for item in parts if _clean_slide_text(item, 175)]
    return [_clean_slide_text(content, 175)] if content else []


def _ppt_safe_image_path(image_path: str) -> str:
    """Converts image formats unsupported by python-pptx, such as WebP, into PNG."""
    try:
        from PIL import Image

        source = Path(image_path)
        with Image.open(source) as image:
            if image.format in {"BMP", "GIF", "JPEG", "PNG", "TIFF", "WMF"}:
                return str(source)
            converted = source.with_suffix(".png")
            image.convert("RGB").save(converted, format="PNG")
            return str(converted)
    except Exception:
        return image_path


def generate_ppt_from_report(
    report: str,
    query: str,
    output_path: str = "research_presentation.pptx",
    chart_paths: Optional[Sequence[str]] = None,
) -> str:
    """Creates a polished 16:9 .pptx presentation from a generated report."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except Exception as exc:
        raise RuntimeError("python-pptx is required for PPT export. Install it with `pip install python-pptx`.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slides = presentation_planner_agent(report, query)
    section_slides = _polish_slide_plan(_section_slides_from_report(report), query, report)
    if len(slides) < min(5, len(section_slides)) and len(report) > 1200:
        slides = section_slides
    chart_paths = list(chart_paths or [])
    chart_index = 0
    deck_title = _derive_deck_title(query, report)

    navy = RGBColor(15, 23, 42)
    slate = RGBColor(51, 65, 85)
    muted = RGBColor(100, 116, 139)
    paper = RGBColor(248, 250, 252)
    white = RGBColor(255, 255, 255)
    indigo = RGBColor(79, 70, 229)
    cyan = RGBColor(14, 165, 233)
    amber = RGBColor(245, 158, 11)

    def set_background(slide, color: RGBColor = paper) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_textbox(
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int,
        color: RGBColor,
        bold: bool = False,
        align: Optional[Any] = None,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        if align:
            paragraph.alignment = align
        return box

    def add_card(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor = white):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = RGBColor(226, 232, 240)
        return card

    def add_brand_frame(slide, index: int, section: str) -> None:
        set_background(slide)
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), Inches(7.5))
        rail.fill.solid()
        rail.fill.fore_color.rgb = navy
        rail.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(0), Inches(0.06), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = indigo
        accent.line.fill.background()

        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(0.42), Inches(2.15), Inches(0.34))
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(238, 242, 255)
        pill.line.fill.background()
        pill_text = pill.text_frame
        pill_text.clear()
        pill_para = pill_text.paragraphs[0]
        pill_para.text = section.upper()[:28]
        pill_para.font.size = Pt(9)
        pill_para.font.bold = True
        pill_para.font.color.rgb = indigo
        pill_para.alignment = PP_ALIGN.CENTER

        footer = add_textbox(slide, f"Research Agentic Chatbot · {index:02d}", 9.6, 7.04, 2.95, 0.25, 9, muted, align=PP_ALIGN.RIGHT)
        footer.text_frame.margin_top = 0

    def add_notes(slide, notes: str) -> None:
        if not notes:
            return
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = notes[:1200]
        except Exception:
            return

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(title_slide, navy)
    top_bar = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = indigo
    top_bar.line.fill.background()
    orb = title_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.7), Inches(3.8), Inches(3.8))
    orb.fill.solid()
    orb.fill.fore_color.rgb = indigo
    orb.line.fill.background()
    orb2 = title_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.35), Inches(4.95), Inches(2.9), Inches(2.9))
    orb2.fill.solid()
    orb2.fill.fore_color.rgb = cyan
    orb2.line.fill.background()
    add_textbox(title_slide, "RESEARCH BRIEFING", 0.9, 1.25, 4.1, 0.35, 11, cyan, bold=True)
    add_textbox(title_slide, deck_title[:120], 0.88, 2.05, 10.4, 1.55, 38, white, bold=True)
    add_textbox(title_slide, f"{len(slides)} slides · generated from: {query[:90]}", 0.92, 3.82, 8.6, 0.55, 15, RGBColor(203, 213, 225))
    add_textbox(title_slide, "Prepared by Research Agentic Chatbot", 0.92, 6.72, 5.6, 0.3, 10, RGBColor(148, 163, 184))
    add_notes(title_slide, f"Deck generated for the request: {query}\n\nUse this opening to frame the audience, objective, and expected decision.")

    for slide_number, slide_plan in enumerate(slides, start=1):
        title = str(slide_plan.get("title", "")).strip()
        if not title or title.lower() == "title":
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        section = str(slide_plan.get("section", "Insight")).strip() or "Insight"
        add_brand_frame(slide, slide_number, section)
        add_textbox(slide, title[:95], 0.72, 0.86, 11.4, 0.72, 27, navy, bold=True)

        wants_chart = slide_plan.get("visual") in {"chart", "image"} or any(
            keyword in title.lower()
            for keyword in ["data", "market", "analysis", "trend", "correlation", "comparison"]
        )
        has_chart = wants_chart and chart_index < len(chart_paths) and os.path.exists(chart_paths[chart_index])
        bullets = _normalize_slide_bullets(slide_plan.get("bullets", []))[:5] or ["See generated report for details."]
        takeaway = str(slide_plan.get("takeaway", "")).strip()
        speaker_notes = str(slide_plan.get("speaker_notes", "")).strip()
        slide_type = str(slide_plan.get("slide_type", "insight")).lower()
        if not speaker_notes:
            speaker_notes = f"{title}\n\n" + "\n".join(f"- {bullet}" for bullet in bullets)

        body_left = 0.82
        body_top = 2.02
        body_width = 7.05 if has_chart else 11.35
        add_card(slide, body_left, body_top, body_width, 3.9)
        body_box = slide.shapes.add_textbox(Inches(body_left + 0.28), Inches(body_top + 0.28), Inches(body_width - 0.56), Inches(3.32))
        body = body_box.text_frame
        body.clear()
        body.word_wrap = True
        body.margin_left = Inches(0.03)
        body.margin_right = Inches(0.03)

        bullet_font_size = 13 if slide_type == "reference" else 15 if len(bullets) >= 5 else 17
        bullet_space = 7 if slide_type == "reference" else 10 if len(bullets) >= 5 else 12

        for bullet_index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if bullet_index == 0 else body.add_paragraph()
            paragraph.text = bullet[:185]
            paragraph.level = 0
            paragraph.font.size = Pt(bullet_font_size)
            paragraph.font.color.rgb = slate
            paragraph.space_after = Pt(bullet_space)
            paragraph.line_spacing = 1.08

        if takeaway and slide_type != "reference":
            callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(6.12), Inches(11.35), Inches(0.62))
            callout.fill.solid()
            callout.fill.fore_color.rgb = RGBColor(255, 251, 235)
            callout.line.color.rgb = RGBColor(253, 230, 138)
            add_textbox(slide, f"Key takeaway: {takeaway[:170]}", 1.05, 6.27, 10.8, 0.27, 12, RGBColor(120, 53, 15), bold=True)
        elif slide_type == "reference":
            add_textbox(slide, "Source list captured from the generated report. Verify important claims before publication.", 0.92, 6.21, 10.95, 0.38, 11, muted)

        if has_chart:
            image_path = _ppt_safe_image_path(chart_paths[chart_index])
            add_card(slide, 8.25, 2.02, 4.35, 3.9)
            slide.shapes.add_picture(image_path, Inches(8.45), Inches(2.3), width=Inches(3.95))
            add_textbox(slide, "Supporting visual", 8.46, 5.55, 3.95, 0.25, 9, muted, align=PP_ALIGN.CENTER)
            chart_index += 1
        elif not takeaway:
            metric = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.85), Inches(5.62), Inches(0.72), Inches(0.72))
            metric.fill.solid()
            metric.fill.fore_color.rgb = amber
            metric.line.fill.background()

        add_notes(slide, speaker_notes)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return str(output)


# ==========================================
# 4. NODE FUNCTIONS
# ==========================================
def input_preprocessor_node(state: ResearchState) -> Dict[str, Any]:
    """Adds PDF, image, data, and math-agent findings to research memory."""
    print("--- 📥 PREPROCESSING RESEARCH INPUTS ---")
    sections = []
    chart_paths: List[str] = []

    for file_path in state.get("source_files", []):
        suffix = Path(file_path).suffix.lower()
        if suffix == ".pdf":
            sections.append(f"\n[PDF INPUT: {file_path}]\n{extract_pdf_text(file_path)}")
        elif suffix in IMAGE_EXTENSIONS:
            sections.append(f"\n[IMAGE INPUT: {file_path}]\n{analyze_image(file_path)}")
        elif suffix in DATA_EXTENSIONS:
            analysis = run_local_data_analysis(file_path, query=state["query"])
            sections.append(f"\n[DATA ANALYSIS INPUT: {file_path}]\n{analysis['summary']}")
            chart_paths.extend(analysis.get("chart_paths", []))
        else:
            sections.append(f"\n[UNSUPPORTED INPUT SKIPPED: {file_path}]")

    math_solution = solve_math_problem(state["query"])
    if math_solution:
        sections.append(f"\n[MATH AGENT OUTPUT]\n{math_solution}")

    input_context = "\n\n".join(sections)
    return {
        "input_context": input_context,
        "research_data": f"\n\n--- User Provided Context & Agent Analysis ---\n{input_context}" if input_context else "",
        "chart_paths": chart_paths,
    }


def planner_node(state: ResearchState) -> Dict[str, Any]:
    """Creates a comprehensive, step-by-step research plan."""
    print("--- 📝 CREATING RESEARCH PLAN ---")
    prompt = f"""
    Create a rigorous deep-research plan to answer this query comprehensively: '{state['query']}'.
    Resolve relative dates against this current time context:
    {build_current_datetime_context()}

    Use this compact memory/context only if relevant:
    {state.get('memory_context', '')}

    Account for this user-provided file/tool context if present:
    {state.get('input_context', '')}

    Requirements:
    - If DATA ANALYSIS INPUT is present, make uploaded dataset analysis the primary evidence and only use external search when the user explicitly asks for external/current context.
    - Search across official sources, primary sources, reputable news/research sites, and high-signal secondary sources.
    - Identify what quantitative data, calculations, comparisons, timelines, prices, market sizes, scores, schedules, or estimates are needed.
    - Plan at least one numerical/analytical pass whenever the topic has measurable facts.
    - Include cross-verification, contradictions to check, citation gathering, limitations, recommendations, and presentation-ready structure.
    - Prefer exact dates, units, currencies, assumptions, and source URLs over vague claims.
    Output only the plan.
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    return {"plan": response.content, "iteration": 0}


def search_and_summarize_node(state: ResearchState) -> Dict[str, Any]:
    """An autonomous search agent that routes queries between Google and Wikipedia."""
    print(f"--- 🔍 INTELLIGENT SEARCH AGENT (Iteration {state['iteration'] + 1}) ---")

    has_uploaded_data = "[DATA ANALYSIS INPUT:" in state.get("input_context", "")
    external_request = bool(re.search(r"\b(web|internet|latest|current|today|source|sources|cite|citation|external|market|news)\b", state["query"], re.IGNORECASE))
    if has_uploaded_data and not external_request:
        return {
            "research_data": (
                f"\n\n--- Iteration {state['iteration'] + 1} Local Dataset Findings ---\n"
                "External web search skipped because uploaded data analysis is available and the user did not ask for external/current sources. "
                "Use the DATA ANALYSIS INPUT, generated charts, descriptive statistics, correlations, outliers, and limitations as the primary evidence."
            ),
            "iteration": state["iteration"] + 1,
        }

    search_target = state["current_gaps"] if state["iteration"] > 0 else state["plan"]
    routing_prompt = f"""
    You are an expert research librarian. Look at the current research target: '{search_target}'.
    Current time context:
    {build_current_datetime_context()}

    Generate up to 6 highly specific search queries to gather this information.

    For each query, choose the best source:
    - "wiki": For established facts, history, core concepts, and mathematical formulas.
    - "web": For recent news, live data, official pages, specific papers, schedules, prices, datasets, and niche implementations.

    Query design rules:
    - Include official-source searches where applicable.
    - Include one query for quantitative data/calculations if the topic has numbers.
    - Include one query for criticism/risks/limitations.
    - Include current year/date terms for time-sensitive topics.

    Output EXACTLY a JSON array of dictionaries, nothing else.
    Format: [{{"source": "wiki", "query": "Transformer neural network"}}, {{"source": "web", "query": "State Space Models latest papers 2024"}}]
    """
    routing_response = llm.invoke([HumanMessage(content=routing_prompt)])
    clean_json = routing_response.content.replace("```json", "").replace("```", "").strip()

    try:
        queries = json.loads(clean_json)
    except json.JSONDecodeError:
        print("⚠️ Agent returned invalid JSON, falling back to a default web search.")
        queries = [{"source": "web", "query": search_target}]

    raw_results = ""
    for query_item in queries:
        source = query_item.get("source")
        query_text = query_item.get("query")
        print(f"    -> Querying {str(source).upper()}: '{query_text}'")
        try:
            if source == "wiki":
                raw_results += f"\n[WIKIPEDIA RESULT FOR '{query_text}']:\n" + wiki_search_tool.run(query_text)
            else:
                raw_results += f"\n[WEB RESULT FOR '{query_text}']:\n" + serper_search(query_text)
        except Exception as exc:
            raw_results += f"\n[Search failed for '{query_text}']: {exc}"

    summary_prompt = f"""
    You are an expert research analyst and data extractor.
    Extract:
    - Key factual claims with source names and URLs.
    - Exact dates, prices, quantities, units, formulas, tables, and comparable data.
    - Any calculations that can be performed from the retrieved numbers; show assumptions and arithmetic.
    - Contradictions, uncertainty, missing data, risks, limitations, and recommendations.
    - Useful code snippets or reproducible analysis steps only if relevant.
    Ignore conversational fluff. Preserve mathematical notation, currency, units, and citation details.

    Raw Results:
    {raw_results}
    """
    summary_response = llm.invoke([HumanMessage(content=summary_prompt)])
    return {
        "research_data": f"\n\n--- Iteration {state['iteration'] + 1} Findings ---\n" + summary_response.content,
        "iteration": state["iteration"] + 1,
    }


def verification_node(state: ResearchState) -> Dict[str, Any]:
    """Cross-verifies claims, checks math/code, and identifies gaps."""
    print("--- ⚖️ VERIFYING DATA & CHECKING GAPS ---")
    prompt = f"""
    You are a rigorous academic reviewer and fact-checker. Evaluate the accumulated research against the original query.
    Original Query: {state['query']}
    Current time context:
    {build_current_datetime_context()}

    Accumulated Research:
    {state['research_data']}

    Tasks:
    1. Cross-verify claims for contradictions and stale information.
    2. Check mathematical calculations, dates, prices, schedules, data-analysis claims, and code snippets for errors.
    3. Identify missing information required to fully answer the query.
    4. Check that citations, source URLs, tables, limitations, assumptions, and recommendations are available where needed.
    5. Demand another search if the answer lacks official/current sources for time-sensitive topics.

    If the data is complete and accurate, output exactly: COMPLETE.
    If there are errors or missing data, output a detailed list of the gaps that need to be researched next.
    """
    response = llm.invoke([HumanMessage(content=prompt)])
    return {"current_gaps": response.content}


def format_agent_node(state: ResearchState) -> Dict[str, Any]:
    """Formats the final report into LaTeX, Markdown, Plain Text, or PPT-ready Markdown."""
    print(f"--- 🖨️ FORMATTING OUTPUT AS {state['format_preference'].upper()} ---")

    output_format = "markdown" if state["format_preference"].lower() == "pptx" else state["format_preference"]
    system_instruction = f"""
    You are a professional deep-research report generator. Create an excellent, polished, evidence-backed report from verified research.
    Resolve all relative dates against this current time context:
    {build_current_datetime_context()}

    Required report tasks:
    - Start with a direct answer and executive summary.
    - Organize findings with clear headings.
    - Include calculations/quantitative analysis where applicable, showing assumptions and arithmetic.
    - Add comparison tables, timelines, price/schedule tables, or evidence matrices where useful.
    - Cite source names and URLs inline for important claims.
    - Add limitations, contradictions, confidence level, recommendations, and references.
    You MUST output the final report entirely in this format: {output_format}

    Rules for LaTeX: provide a full compilable document and wrap math correctly.
    Rules for Markdown: use heading hierarchies, tables, citations/references, and code blocks where useful.
    Rules for Plain Text: use spacing and capitalization to create structure without markdown symbols.
    """
    user_prompt = f"Original Query: {state['query']}\n\nVerified Research Data:\n{state['research_data']}\n\nWrite the final report."
    if state.get("memory_context"):
        user_prompt = (
            f"Original Query: {state['query']}\n\n"
            f"Relevant Memory Context:\n{state['memory_context']}\n\n"
            f"Verified Research Data:\n{state['research_data']}\n\n"
            "Write the final report."
        )
    response = llm.invoke([SystemMessage(content=system_instruction), HumanMessage(content=user_prompt)])
    return {"final_report": response.content}


def ppt_export_node(state: ResearchState) -> Dict[str, Any]:
    """Exports a generated report to PowerPoint when requested."""
    if state["format_preference"].lower() != "pptx":
        return {}

    print("--- 📊 EXPORTING POWERPOINT ---")
    pptx_path = state.get("pptx_path") or "research_presentation.pptx"
    visual_paths = list(state.get("chart_paths", []))
    if len(visual_paths) < 3:
        media_dir = Path(pptx_path).parent
        searched_images = collect_slide_images(
            state["query"],
            output_dir=media_dir,
            plan=state.get("plan"),
            max_images=4 - len(visual_paths),
        )
        visual_paths.extend(searched_images)
    exported_path = generate_ppt_from_report(
        report=state["final_report"],
        query=state["query"],
        output_path=pptx_path,
        chart_paths=visual_paths,
    )
    return {"pptx_path": exported_path, "final_report": f"PowerPoint generated: {exported_path}\n\nSource report:\n\n{state['final_report']}"}


# ==========================================
# 5. CONDITIONAL ROUTING LOGIC
# ==========================================
def should_continue_research(state: ResearchState) -> str:
    """Decides whether to loop back to search or proceed to formatting."""
    has_uploaded_data = "[DATA ANALYSIS INPUT:" in state.get("input_context", "")
    external_request = bool(re.search(r"\b(web|internet|latest|current|today|source|sources|cite|citation|external|market|news)\b", state["query"], re.IGNORECASE))
    if has_uploaded_data and not external_request and state["iteration"] >= 1:
        print("--- 📊 LOCAL DATA ANALYSIS COMPLETE: Proceeding to format ---")
        return "format"

    if state["iteration"] >= 4:
        print("--- 🛑 ITERATION CAP REACHED: Proceeding to format ---")
        return "format"

    if "COMPLETE" in state["current_gaps"]:
        print("--- ✅ VERIFICATION PASSED: Proceeding to format ---")
        return "format"

    print("--- ⚠️ GAPS FOUND: Looping back to search ---")
    return "research"


# ==========================================
# 6. GRAPH COMPILATION & ENCAPSULATION
# ==========================================
def create_research_engine():
    """Builds and compiles the LangGraph application."""
    workflow = StateGraph(ResearchState)

    workflow.add_node("input_preprocessor", input_preprocessor_node)
    workflow.add_node("planner", planner_node)
    workflow.add_node("searcher", search_and_summarize_node)
    workflow.add_node("verifier", verification_node)
    workflow.add_node("formatter", format_agent_node)
    workflow.add_node("ppt_exporter", ppt_export_node)

    workflow.add_edge(START, "input_preprocessor")
    workflow.add_edge("input_preprocessor", "planner")
    workflow.add_edge("planner", "searcher")
    workflow.add_edge("searcher", "verifier")
    workflow.add_conditional_edges("verifier", should_continue_research, {"research": "searcher", "format": "formatter"})
    workflow.add_edge("formatter", "ppt_exporter")
    workflow.add_edge("ppt_exporter", END)

    return workflow.compile()


# =========================================
# 7. MASTER EXECUTION FUNCTION
# ==========================================
def run_advanced_research(
    query: str,
    format_type: str = "markdown",
    pdf_paths: Optional[Union[str, Sequence[str]]] = None,
    image_paths: Optional[Union[str, Sequence[str]]] = None,
    data_paths: Optional[Union[str, Sequence[str]]] = None,
    source_files: Optional[Union[str, Sequence[str]]] = None,
    pptx_path: str = "research_presentation.pptx",
    memory_store: Optional[MemoryStore] = None,
    memory_records: Optional[Sequence[MemoryRecord]] = None,
    user_profile: Optional[Dict[str, Any]] = None,
    conversation_summary: str = "",
    recent_messages: Optional[Sequence[str]] = None,
) -> str:
    """
    Entry point to run the advanced research engine.

    format_type options: 'latex', 'markdown', 'plain text', 'pptx'
    Inputs: PDF paths, image paths, CSV/Excel data paths, mixed source_files,
    and compact memory inputs. Full memory should stay outside the prompt.
    """
    engine = create_research_engine()
    all_source_files = _as_list(source_files) + _as_list(pdf_paths) + _as_list(image_paths) + _as_list(data_paths)
    active_memory_store = memory_store or InMemoryMemoryStore(memory_records)
    context_builder = ContextBuilder(memory_store=active_memory_store)
    built_context = context_builder.build(
        query=query,
        format_type=format_type,
        user_profile=user_profile,
        conversation_summary=conversation_summary,
        recent_messages=recent_messages,
    )

    initial_state = {
        "query": query,
        "format_preference": format_type,
        "memory_context": built_context["context"],
        "task_intent": built_context["intent"],
        "input_context": "",
        "source_files": all_source_files,
        "chart_paths": [],
        "pptx_path": pptx_path,
        "plan": "",
        "research_data": "",
        "current_gaps": "",
        "iteration": 0,
        "final_report": "",
    }

    print(f"\n🚀 Starting Advanced Research for: '{query}'\n")
    final_state = engine.invoke(initial_state)
    return final_state["final_report"]


# --- Example Usage ---
# if __name__ == "__main__":
#     report = run_advanced_research(
#         query="Analyze the market outlook from this PDF and dataset.",
#         format_type="pptx",
#         pdf_paths=["report.pdf"],
#         image_paths=["chart.png"],
#         data_paths=["market.csv"],
#         pptx_path="outputs/research_presentation.pptx",
#     )
#     print(report)
